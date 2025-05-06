from langchain.text_splitter import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
import fitz  # PyMuPDF
import os
import faiss
import numpy as np
import pickle
from docx import Document
from pptx import Presentation
import pandas as pd

UPLOAD_FOLDER_FILES = "uploads/files"
UPLOAD_FOLDER_RAG = "uploads/rag"
ALLOWED_EXTENSIONS = {'.pdf', '.txt', '.docx', '.pptx', '.xls', '.xlsx'}

def extract_text_from_pdf(filepath):
    text = ""
    try:
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        print(f"❌ Failed to extract PDF: {filepath} — {e}")
    return text

def extract_text_from_docx(filepath):
    try:
        doc = Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        print(f"❌ Failed to extract DOCX: {filepath} — {e}")
        return ""

def extract_text_from_pptx(filepath):
    text = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"❌ Failed to extract PPTX: {filepath} — {e}")
    return text

def extract_text_from_excel(filepath):
    text = ""
    try:
        # Try all sheets, all cells
        excel_data = pd.read_excel(filepath, sheet_name=None, engine='openpyxl' if filepath.endswith('.xlsx') else 'xlrd')
        for sheet_name, sheet_df in excel_data.items():
            text += f"Sheet: {sheet_name}\n"
            text += sheet_df.fillna("").astype(str).to_string(index=False, header=False)
            text += "\n\n"
    except Exception as e:
        print(f"❌ Failed to extract Excel: {filepath} — {e}")
    return text

def collect_text_from_folder(folder):
    combined_text = ""
    processed_files = []

    if not os.path.exists(folder):
        print(f"⚠️ Folder does not exist: {folder}")
        return combined_text, processed_files

    for filename in os.listdir(folder):
        filepath = os.path.join(folder, filename)
        ext = os.path.splitext(filename)[1].lower()
        if not os.path.isfile(filepath) or ext not in ALLOWED_EXTENSIONS:
            continue

        print(f"🔍 Processing {filename}")
        processed_files.append(filename)

        try:
            if ext == ".pdf":
                combined_text += extract_text_from_pdf(filepath) + "\n"
            elif ext == ".txt":
                with open(filepath, "r", encoding="utf-8") as f:
                    combined_text += f.read() + "\n"
            elif ext == ".docx":
                combined_text += extract_text_from_docx(filepath) + "\n"
            elif ext == ".pptx":
                combined_text += extract_text_from_pptx(filepath) + "\n"
            elif ext in {".xls", ".xlsx"}:
                combined_text += extract_text_from_excel(filepath) + "\n"
        except Exception as e:
            print(f"⚠️ Failed to process {filename}: {e}")

    return combined_text, processed_files


# Aggregate all RAG-eligible content
text1, files1 = collect_text_from_folder(UPLOAD_FOLDER_FILES)
text2, files2 = collect_text_from_folder(UPLOAD_FOLDER_RAG)
all_text = text1 + text2
all_files = files1 + files2

print(f"📄 Indexed {len(all_files)} files: {all_files}")

print(f"🗂️ Found {len(all_text.split())} words of raw text.")

# Chunk and embed
splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
chunks = splitter.split_text(all_text)

print(f"🧩 Split into {len(chunks)} chunks.")

if not chunks:
    print("❌ No valid chunks to embed. Exiting.")
    exit(1)

model = SentenceTransformer("all-MiniLM-L6-v2")
embeddings = model.encode(chunks)

index = faiss.IndexFlatL2(embeddings[0].shape[0])
index.add(np.array(embeddings))

faiss.write_index(index, "vector_index.faiss")
with open("chunks.pkl", "wb") as f:
    pickle.dump(chunks, f)

print(f"✅ Rebuilt FAISS index with {len(chunks)} chunks.")
